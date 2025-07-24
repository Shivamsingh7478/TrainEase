# extract_images.py
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images(pptx_path, output_dir):
    print(f"[DEBUG] PPTX path received: {pptx_path}")
    print(f"[DEBUG] Output directory: {output_dir}")

    if not os.path.exists(pptx_path):
        print(f"[ERROR] File does not exist at: {pptx_path}")
        return

    if not os.path.exists(output_dir):
        print(f"[DEBUG] Creating output directory at {output_dir}")
        os.makedirs(output_dir)
    else:
        print(f"[DEBUG] Output directory already exists.")

    try:
        prs = Presentation(pptx_path)
        print(f"[DEBUG] Presentation loaded successfully.")
    except Exception as e:
        print(f"[ERROR] Failed to load presentation: {e}")
        return

    total_slides = len(prs.slides)
    print(f"[DEBUG] Total slides found: {total_slides}")
    
    image_count = 0

    for i, slide in enumerate(prs.slides, start=1):
        print(f"[DEBUG] Processing slide {i}...")
        for j, shape in enumerate(slide.shapes, start=1):
            print(f"  [DEBUG] Checking shape {j} of type: {shape.shape_type}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"    [DEBUG] Image found in slide {i}, shape {j}")
                try:
                    image = shape.image
                    ext = image.ext
                    image_bytes = image.blob
                    filename = f"slide{i}_img{j}.{ext}"
                    filepath = os.path.join(output_dir, filename)
                    print(f"    [DEBUG] Saving image to {filepath}")
                    with open(filepath, "wb") as f:
                        f.write(image_bytes)
                    image_count += 1
                except Exception as e:
                    print(f"    [ERROR] Failed to save image: {e}")
            else:
                print(f"    [DEBUG] Not an image, skipping...")

    print(f"[INFO] Total images extracted: {image_count}")
    if image_count == 0:
        print("[WARN] No images were found in the presentation.")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("[USAGE] python extract_images.py <path_to_pptx> <output_folder>")
    else:
        extract_images(sys.argv[1], sys.argv[2])
